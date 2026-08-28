"""Renders team_presentation.pptx to PNG previews for visual QA.

No LibreOffice available in this environment, so this walks the pptx shape
tree with python-pptx and redraws it with PIL at a fixed DPI. It is a QA
approximation (font metrics differ from real PowerPoint), not a renderer --
good enough to catch overflow, overlap and placement mistakes before handoff.
"""

from __future__ import annotations

import io
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

HERE = Path(__file__).parent
DECK = HERE / "team_presentation.pptx"
OUT = HERE / "preview"
OUT.mkdir(exist_ok=True)

DPI = 100  # pixels per inch
FONTDIR = Path(
    __import__("matplotlib").get_data_path()
) / "fonts/ttf"

FONT_CACHE = {}


def font(name, size_pt):
    key = (name, size_pt)
    if key not in FONT_CACHE:
        px = int(size_pt * DPI / 72)
        FONT_CACHE[key] = ImageFont.truetype(str(FONTDIR / name), px)
    return FONT_CACHE[key]


def emu_to_px(emu):
    return int(Emu(emu) / 914400 * DPI)


# defaults pulled directly from agency-template.pptx's slideMaster2.xml
# titleStyle/bodyStyle (p:titleStyle sz=4000 b=1 Arial tx1;
# p:bodyStyle lvl1 sz=2800 Arial tx1, bullet char "*" in cyan) -- used when a
# placeholder run has no explicit override, i.e. it's inheriting from the
# template rather than something this script set directly
TITLE_DEFAULT = dict(size=40, bold=True, color="#000000")
# slideLayout7.xml (the Title Slide layout) overrides ctrTitle to 36pt with
# normAutofit -- distinct from the master's 40pt used by Title Only/Content
CENTER_TITLE_DEFAULT = dict(size=36, bold=True, color="#000000")
BODY_DEFAULT = dict(size=28, bold=False, color="#000000")
SUBTITLE_DEFAULT = dict(size=24, bold=False, color="#000000")


def placeholder_default(shape):
    if not shape.is_placeholder:
        return None
    ph_type = shape.placeholder_format.type
    if ph_type == 3:  # CENTER_TITLE
        return CENTER_TITLE_DEFAULT
    if ph_type == 1:  # TITLE
        return TITLE_DEFAULT
    if ph_type == 4:  # SUBTITLE
        return SUBTITLE_DEFAULT
    if ph_type in (2, 7):  # BODY, OBJECT
        return BODY_DEFAULT
    return None


def wrap_text(draw, text, fnt, max_width_px):
    """Greedy word-wrap so the preview reflects PowerPoint's actual
    word_wrap=True behaviour instead of drawing one unbroken line per
    paragraph -- without this, a too-long title looks identical to a
    correctly-sized one that merely runs off the edge of a crude preview."""
    words = text.split(" ")
    lines, cur = [], ""
    for word in words:
        trial = f"{cur} {word}".strip()
        w = draw.textbbox((0, 0), trial, font=fnt)[2]
        if w <= max_width_px or not cur:
            cur = trial
        else:
            lines.append(cur)
            cur = word
    if cur:
        lines.append(cur)
    return lines or [""]


def draw_deck():
    prs = Presentation(DECK)
    w = emu_to_px(prs.slide_width)
    h = emu_to_px(prs.slide_height)

    for i, slide in enumerate(prs.slides, start=1):
        img = Image.new("RGB", (w, h), "white")
        draw = ImageDraw.Draw(img)

        # decorative pictures (logo, partner-logo strip) live on the layout,
        # not the slide, and are inherited automatically by real PowerPoint --
        # draw them here too so Title Slide previews aren't missing branding
        for shape in slide.slide_layout.shapes:
            if shape.is_placeholder or shape.shape_type != 13:
                continue
            left, top = emu_to_px(shape.left), emu_to_px(shape.top)
            width, height = emu_to_px(shape.width), emu_to_px(shape.height)
            try:
                im = Image.open(io.BytesIO(shape.image.blob)).convert("RGBA")
                im = im.resize((max(width, 1), max(height, 1)))
                bg = Image.new("RGBA", im.size, (255, 255, 255, 255))
                im = Image.alpha_composite(bg, im).convert("RGB")
                img.paste(im, (left, top))
            except Exception:
                pass

        for shape in slide.shapes:
            left, top = emu_to_px(shape.left), emu_to_px(shape.top)
            width, height = emu_to_px(shape.width), emu_to_px(shape.height)

            if shape.shape_type == 13:  # PICTURE
                try:
                    im = Image.open(io.BytesIO(shape.image.blob))
                    im = im.convert("RGB").resize((max(width, 1), max(height, 1)))
                    img.paste(im, (left, top))
                except Exception as e:
                    draw.rectangle([left, top, left + width, top + height], outline="red")
                    draw.text((left, top), f"[image error: {e}]", fill="red")
                continue

            if shape.has_text_frame:
                tf = shape.text_frame
                # rectangles used as rules/cards have text frames too but we
                # only care about fill for those with no text
                has_text = any(p.text.strip() for p in tf.paragraphs)
                if shape.shape_type == 1 and not has_text:  # AUTO_SHAPE / rectangle
                    fill_rgb = None
                    try:
                        fill_rgb = shape.fill.fore_color.rgb
                    except Exception:
                        pass
                    color = f"#{fill_rgb}" if fill_rgb else "#cccccc"
                    draw.rectangle([left, top, left + max(width, 2), top + max(height, 2)], fill=color)
                    continue

                ph_default = placeholder_default(shape)
                is_body_ph = shape.is_placeholder and shape.placeholder_format.type in (2, 7)
                anchor = tf.vertical_anchor

                # Pass 1: resolve each paragraph's formatting, then wrap it
                # into the lines PowerPoint's word_wrap=True would actually
                # produce, so overflow/centering is judged on real line
                # counts rather than one (possibly very long) line per
                # paragraph.
                rows = []  # (text, fnt, color, align, extra_space_after_px)
                for p in tf.paragraphs:
                    text = "".join(r.text for r in p.runs)
                    if is_body_ph and text.strip():
                        text = "•  " + text
                    sz = (ph_default or {}).get("size", 18)
                    bold = (ph_default or {}).get("bold", False)
                    italic = False
                    color = (ph_default or {}).get("color", "#2f2f2f")
                    fam = ""
                    for r in p.runs:
                        if r.font.size:
                            sz = r.font.size.pt
                        if r.font.bold is not None:
                            bold = bool(r.font.bold)
                        italic = bool(r.font.italic)
                        try:
                            if r.font.color.rgb:
                                color = f"#{r.font.color.rgb}"
                        except Exception:
                            pass
                        fam = (r.font.name or "").lower()
                    if "georgia" in fam or "serif" in fam:
                        fname = "DejaVuSerif-BoldItalic.ttf" if (bold and italic) else (
                            "DejaVuSerif-Bold.ttf" if bold else (
                                "DejaVuSerif-Italic.ttf" if italic else "DejaVuSerif.ttf"))
                    else:
                        fname = "DejaVuSans-BoldOblique.ttf" if (bold and italic) else (
                            "DejaVuSans-Bold.ttf" if bold else (
                                "DejaVuSans-Oblique.ttf" if italic else "DejaVuSans.ttf"))
                    fnt = font(fname, sz)
                    align = p.alignment
                    space_after_px = p.space_after.pt * DPI / 72 if p.space_after else 0
                    if not text.strip():
                        rows.append(("", fnt, color, align, sz * DPI / 72 * 1.25))
                        continue
                    wrapped = wrap_text(draw, text, fnt, width)
                    for j, line in enumerate(wrapped):
                        extra = space_after_px if j == len(wrapped) - 1 else 0
                        rows.append((line, fnt, color, align, sz * DPI / 72 * 1.25 + extra))

                total_h = sum(r[4] for r in rows)
                cursor_y = top + max(0, (height - total_h) / 2) if anchor == 3 else top

                for text, fnt, color, align, lh in rows:
                    if text:
                        bbox = draw.textbbox((0, 0), text, font=fnt)
                        tw = bbox[2] - bbox[0]
                        if align == 2:  # CENTER
                            x = left + (width - tw) / 2
                        elif align == 3:  # RIGHT
                            x = left + width - tw
                        else:
                            x = left
                        draw.text((x, cursor_y), text, font=fnt, fill=color)
                    cursor_y += lh

        outpath = OUT / f"slide_{i:02d}.png"
        img.save(outpath)
        print(outpath)


if __name__ == "__main__":
    draw_deck()
