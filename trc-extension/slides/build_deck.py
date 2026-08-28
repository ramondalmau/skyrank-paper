"""Builds the 30-minute team deck for the route-preference model.

Built on the real agency-template.pptx, not a from-scratch design: slides use
the template's actual layouts (Title Slide / Title and Content / Title Only /
Blank) from its branded master (slideMaster2.xml, theme2 -- confirmed by
direct inspection to be the master the template's own Title Slide and closing
"Thank you!" example slides both use), so titles, body bullets, background,
and the corner brand mark all come from the template's own inheritance chain
rather than being reproduced by hand. The one font used for anything built
here is Arial, matching the template's own title/body styles exactly.

Design rule: one idea per slide, almost always a real chart or figure. No
slide's entire content is a single typeset number -- every figure here is
either the paper's own or a small chart built from the paper's own tables.
Chart/figure colour stays in the paper's own palette (paper/sid2026/style.py)
so reused and newly-built charts read as one set; only slide chrome (cards,
accent bars) draws on the template's theme accent colours.
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).parent
FIG = HERE / "figures"
TEMPLATE = HERE / "agency-template.pptx"
OUT = HERE / "team_presentation.pptx"

# paper palette (paper/sid2026/style.py) -- used for chart/figure content
INK = RGBColor(0x2F, 0x2F, 0x2F)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
RED = RGBColor(0xB5, 0x48, 0x4D)
CAPTION_GREY = RGBColor(0x6E, 0x6E, 0x6E)
CARD_BG = RGBColor(0xF6, 0xF5, 0xF3)
BORDER = RGBColor(0xCB, 0xC8, 0xC2)

# template theme accents (theme2.xml) -- used only for slide chrome we add
THEME_NAVY = RGBColor(0x1A, 0x35, 0x65)
THEME_CYAN = RGBColor(0x00, 0x9C, 0xD7)
THEME_GREEN = RGBColor(0x8B, 0xB5, 0x3A)

FONT = "Arial"  # matches the template's own title (Arial Bold) / body (Arial)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# safe content zone: the master's corner brand mark sits at
# (12.57, 6.76, 0.57x0.58) -- i.e. only the bottom-RIGHT corner is reserved.
# Everything we build stays left of x=11.98 (CONTENT_LEFT + CONTENT_WIDTH),
# clear of the mark by 0.59in, so it's safe to use the fuller height down to
# SAFE_BOTTOM rather than stopping at the template's own Content Placeholder
# bottom (6.76), which was sized for 28pt bullet text, not 26pt+ captions.
CONTENT_LEFT = Inches(0.48)
CONTENT_WIDTH = Inches(11.50)
SAFE_BOTTOM = Inches(7.25)

BODY_SIZE = 26  # floor for every piece of text this script sets explicitly


def load_template() -> Presentation:
    prs = Presentation(TEMPLATE)
    # the file ships 5 placeholder/example slides (title, 3 empty "Title and
    # Content", a closing "Thank you!"). Removing them from the slide-id
    # list alone leaves their parts (slide1.xml..slide5.xml) reachable via
    # presentation.xml.rels, so the next add_slide() collides with those
    # partnames -- drop the relationship too so the parts are actually freed.
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides):
        prs.part.drop_rel(sld.rId)
        xml_slides.remove(sld)
    return prs


# slideMaster2.xml / theme2: the master the template's own Title Slide and
# closing "Thank you!" example both use, and the only master with a complete,
# logo-bearing set of layouts confirmed by direct inspection of the pptx XML
BRAND_MASTER = 1


def layouts(prs):
    master = prs.slide_masters[BRAND_MASTER]
    by_name = {layout.name: layout for layout in master.slide_layouts}
    return {
        "title": by_name["Title Slide"],
        "content": by_name["Title and Content"],
        "title_only": by_name["Title Only"],
        "blank": by_name["Blank"],
    }


def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    size,
    color,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    anchor=None,
    line_spacing=None,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = FONT
    return tb


def fit_image(slide, img_path, box_left, box_top, box_w, box_h):
    im = Image.open(img_path)
    iw, ih = im.size
    box_ratio = box_w / box_h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        w = box_w
        h = Emu(int(box_w / img_ratio))
    else:
        h = box_h
        w = Emu(int(box_h * img_ratio))
    left = Emu(int(box_left + (box_w - w) / 2))
    top = Emu(int(box_top + (box_h - h) / 2))
    slide.shapes.add_picture(str(img_path), left, top, width=w, height=h)


def add_rule(slide, left, top, width, color, weight_pt=4):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(weight_pt))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.shadow.inherit = False
    return line


# --------------------------------------------------------------------------
# slide builders -- each uses a real template layout + its real placeholders
# for anything the template defines (title, background, bullets); only the
# picture, caption, and card/box chrome are ours
# --------------------------------------------------------------------------


def set_placeholder_text(placeholder, text, size=BODY_SIZE, height=None):
    """Placeholders inherit their font size from the layout/master (e.g. the
    Title Slide layout's Subtitle is 24pt in a box sized for that) -- override
    explicitly so every piece of text meets the size floor, and grow the box
    to match rather than letting the extra lines overflow it.

    A placeholder with no local <a:xfrm> reports inherited left/top/width/
    height when read, but python-pptx writes a fresh xfrm with only the one
    attribute touched -- setting height alone silently zeroes width. Read all
    four first (resolving inheritance) and write all four back explicitly."""
    left, top, width = placeholder.left, placeholder.top, placeholder.width
    cur_height = placeholder.height
    placeholder.text = text
    for run in placeholder.text_frame.paragraphs[0].runs:
        run.font.size = Pt(size)
    placeholder.left = left
    placeholder.top = top
    placeholder.width = width
    placeholder.height = height if height is not None else cur_height


def title_slide(prs, L, title, subtitle, footer=None):
    slide = prs.slides.add_slide(L["title"])
    slide.shapes.title.text = title
    # subtitle box is 6.41in wide by design (left half is the logo); at 26pt
    # that wraps to 3-4 lines, so it needs more height than the template's
    # own 0.87in -- there's clear space below it until the partner-logo
    # strip at y=6.30, so grow into that room
    set_placeholder_text(slide.placeholders[1], subtitle, height=Inches(1.8))
    if footer:
        # kept left of the partner-logo strip (which starts at x=8.10in)
        # and low enough to clear the subtitle placeholder above it
        add_text(slide, footer, Inches(0.34), Inches(5.85), Inches(7.5), Inches(1.3),
                  BODY_SIZE, CAPTION_GREY)
    return slide


def image_slide(prs, L, title, img_name, caption):
    slide = prs.slides.add_slide(L["title_only"])
    slide.shapes.title.text = title
    caption_h = Inches(1.3)
    caption_top = Emu(int(SAFE_BOTTOM) - int(caption_h))
    image_h = Emu(int(caption_top) - Inches(0.15) - Inches(1.85))
    fit_image(slide, FIG / img_name, CONTENT_LEFT, Inches(1.85), CONTENT_WIDTH, image_h)
    add_text(slide, caption, CONTENT_LEFT, caption_top, CONTENT_WIDTH, caption_h,
              BODY_SIZE, CAPTION_GREY, italic=True, align=PP_ALIGN.CENTER)
    return slide


def limits_slide(prs, L, title, items):
    slide = prs.slides.add_slide(L["content"])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = items[0]
    for item in items[1:]:
        p = tf.add_paragraph()
        p.text = item
    # the template's own bodyStyle is 28pt, already above the 26pt floor --
    # left as the template sets it rather than overridden
    return slide


def flow_slide(prs, L, title, steps, caption):
    """A vertical numbered sequence -- chosen over side-by-side boxes because
    at a 26pt+ floor, boxes narrow enough to sit 4-5 across can't hold a full
    sentence without cramming; a full-width row per step can."""
    slide = prs.slides.add_slide(L["title_only"])
    slide.shapes.title.text = title
    n = len(steps)
    top0 = Inches(1.9)
    caption_h = Inches(0.9)
    caption_top = Emu(int(SAFE_BOTTOM) - int(caption_h))
    steps_bottom = Emu(int(caption_top) - Inches(0.15))
    row_h = Emu((int(steps_bottom) - int(top0)) // n)
    num_w = Inches(0.75)
    for i, (heading, detail) in enumerate(steps):
        row_top = Emu(int(top0) + i * int(row_h))
        tb = slide.shapes.add_textbox(CONTENT_LEFT, row_top,
                                        Emu(int(CONTENT_WIDTH)), row_h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r_num = p.add_run()
        r_num.text = f"{i + 1:02d}   "
        r_num.font.size = Pt(BODY_SIZE)
        r_num.font.bold = True
        r_num.font.color.rgb = THEME_CYAN
        r_num.font.name = FONT
        r_head = p.add_run()
        r_head.text = f"{heading} — "
        r_head.font.size = Pt(BODY_SIZE)
        r_head.font.bold = True
        r_head.font.color.rgb = INK
        r_head.font.name = FONT
        r_detail = p.add_run()
        r_detail.text = detail
        r_detail.font.size = Pt(BODY_SIZE)
        r_detail.font.bold = False
        r_detail.font.color.rgb = INK
        r_detail.font.name = FONT
        if i < n - 1:
            add_rule(slide, CONTENT_LEFT, Emu(int(row_top) + int(row_h) - Inches(0.02)),
                      CONTENT_WIDTH, BORDER, weight_pt=1)
    add_text(slide, caption, CONTENT_LEFT, caption_top, CONTENT_WIDTH, caption_h,
              BODY_SIZE, CAPTION_GREY, italic=True, align=PP_ALIGN.CENTER)
    return slide


def two_card_slide(prs, L, title, cards, caption=None, stacked=False):
    """label gets a 1-line-ish budget (up to 2 short lines); desc gets the
    rest of the card down to a fixed bottom margin -- sized generously
    since descriptions here run 3 lines (count / structure / what it trains).

    stacked=True swaps the two side-by-side (half-width) cards for two
    full-width cards, one above the other. Half-width cards give only
    ~4.65in of usable text -- ~23 characters before wrapping at 26pt --
    which fits a short tag ("Which route wins?") but not an actual
    definition sentence. Full width roughly doubles that budget."""
    slide = prs.slides.add_slide(L["title_only"])
    slide.shapes.title.text = title
    if stacked:
        n = len(cards)
        card_w = CONTENT_WIDTH
        top0 = Inches(1.65)
        gap = Inches(0.15)
        caption_h = Inches(0.85) if caption else Inches(0)
        caption_top = Emu(int(SAFE_BOTTOM) - int(caption_h))
        cards_bottom = Emu(int(caption_top) - (Inches(0.15) if caption else 0))
        card_h = Emu((int(cards_bottom) - int(top0) - int(gap) * (n - 1)) // n)
        label_h = Inches(0.5)
        desc_top_offset = Inches(0.78)
        desc_h = Emu(int(card_h) - int(desc_top_offset) - Inches(0.08))
        for i, (label, desc, color) in enumerate(cards):
            top = Emu(int(top0) + i * (int(card_h) + int(gap)))
            card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, CONTENT_LEFT, top, card_w, card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = CARD_BG
            card.line.color.rgb = BORDER
            card.line.width = Pt(1)
            card.shadow.inherit = False
            add_rule(slide, Emu(int(CONTENT_LEFT) + Inches(0.35)), Emu(int(top) + Inches(0.2)), Inches(0.9), color)
            add_text(slide, label, Emu(int(CONTENT_LEFT) + Inches(0.35)), Emu(int(top) + Inches(0.26)),
                      Emu(int(card_w) - Inches(0.7)), label_h, BODY_SIZE, INK, bold=True)
            add_text(slide, desc, Emu(int(CONTENT_LEFT) + Inches(0.35)), Emu(int(top) + int(desc_top_offset)),
                      Emu(int(card_w) - Inches(0.7)), desc_h, BODY_SIZE, CAPTION_GREY)
        if caption:
            add_text(slide, caption, CONTENT_LEFT, caption_top, CONTENT_WIDTH, caption_h,
                      BODY_SIZE, CAPTION_GREY, italic=True, align=PP_ALIGN.CENTER)
        return slide
    card_w = Inches(5.35)
    card_h = Inches(3.8)
    label_h = Inches(1.3)
    desc_top_offset = Inches(0.65) + label_h + Inches(0.1)
    desc_h = Emu(int(card_h) - int(desc_top_offset) - Inches(0.2))
    gap = Inches(0.5)
    total_w = Emu(int(card_w) * 2 + int(gap))
    left0 = Emu(int((int(SLIDE_W) - int(total_w)) / 2))
    top = Inches(2.1)
    for i, (label, desc, color) in enumerate(cards):
        left = Emu(int(left0) + i * (int(card_w) + int(gap)))
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER
        card.line.width = Pt(1)
        card.shadow.inherit = False
        add_rule(slide, Emu(int(left) + Inches(0.35)), Emu(int(top) + Inches(0.35)), Inches(0.9), color)
        add_text(slide, label, Emu(int(left) + Inches(0.35)), Emu(int(top) + Inches(0.65)),
                  Emu(int(card_w) - Inches(0.7)), label_h, BODY_SIZE, INK, bold=True)
        add_text(slide, desc, Emu(int(left) + Inches(0.35)), Emu(int(top) + int(desc_top_offset)),
                  Emu(int(card_w) - Inches(0.7)), desc_h, BODY_SIZE, CAPTION_GREY)
    if caption:
        add_text(slide, caption, CONTENT_LEFT, Emu(int(top) + int(card_h) + Inches(0.25)),
                  CONTENT_WIDTH, Inches(1.0), BODY_SIZE, CAPTION_GREY, italic=True,
                  align=PP_ALIGN.CENTER)
    return slide


def closing_slide(prs, L, title, subtitle):
    slide = prs.slides.add_slide(L["title"])
    slide.shapes.title.text = title
    set_placeholder_text(slide.placeholders[1], subtitle, height=Inches(1.2))
    return slide


# --------------------------------------------------------------------------
# the deck
# --------------------------------------------------------------------------

def build():
    prs = load_template()
    L = layouts(prs)

    title_slide(
        prs, L,
        "Learning what airlines actually prefer",
        "Every flight-plan revision records a choice the airline made — this\n"
        "is what those choices reveal",
        "EUROCONTROL Innovation Hub  ·  Maastricht Upper Area Control Centre",
    )

    # --- what we're building, and the core trick ---

    flow_slide(
        prs, L,
        "The goal: proposals airlines will actually accept",
        [
            ("A cheaper route is not enough",
             "If the airline would not have chosen it, the proposal is ignored"),
            ("The system must learn preference",
             "Not just minimise fuel or time, but predict which route this airline would accept"),
            ("Preference is learned, not assumed",
             "From what airlines have actually chosen before, not from a cost formula"),
        ],
        "An ignored proposal channel is a channel that eventually gets switched off.",
    )

    image_slide(
        prs, L,
        "Pairs, not sets",
        "pair_idea.png",
        "Earlier work compared the filed route to a SET of generated alternatives. "
        "This work compares PAIRS from real revisions instead.",
    )

    image_slide(
        prs, L,
        "Every revision is a trade",
        "map_escape.png",
        "Amsterdam → Barcelona: 39 extra minutes and 1.4 tonnes of fuel "
        "to avoid 102 minutes of network delay.",
    )

    image_slide(
        prs, L,
        "Not every revision is about delay",
        "map_saver.png",
        "No regulation involved — the operator identified a cheaper "
        "route and filed it directly.",
    )

    image_slide(
        prs, L,
        "What the archive records",
        "fig_composition.png",
        "Most route changes are routine — not the dramatic cases just shown.",
    )

    # --- cost rules, and why they fail ---

    flow_slide(
        prs, L,
        "How a cost rule is tested",
        [
            ("Take a held-out pair", "Two real routes from one revision, never used to train anything"),
            ("Apply the rule", "For example, “prefer less fuel” — the rule always picks one of the two routes"),
            ("Compare to what happened", "Did the rule's pick match the route the airline actually filed?"),
            ("Accuracy = share correct", "Across all such pairs, “prefer less fuel” is right only 44.9% of the time"),
        ],
        "The same test, repeated for every rule in the next chart — only one beats a coin flip.",
    )

    image_slide(
        prs, L,
        "Cost rules score below chance",
        "fig_heuristics.png",
        "Simple cost rules all underperform a coin flip. Only “avoid delay” helps.",
    )

    image_slide(
        prs, L,
        "How it works, start to finish",
        "fig_method.png",
        "Revisions and stay pairs both become training examples. The "
        "model learns the pattern and estimates its own confidence.",
    )

    # --- the bias correction: two datasets, why both, and how the test works ---

    two_card_slide(
        prs, L,
        "Two datasets, two questions",
        [
            ("Revision pairs  —  884,000 pairs",
             "One flight refiles: it files route A, then switches to route B "
             "before departure. The pair is (A, B) — which did the airline pick?", THEME_GREEN),
            ("Stay pairs  —  6,672 pairs",
             "Two flights, same day, same route choice. One was hit by a "
             "regulation, had a real alternative (its twin's route), and stayed. "
             "Pair: (kept, declined).", THEME_CYAN),
        ],
        "Both feed one model — stay pairs count 10× each to offset "
        "being 132× rarer, then held out to test the fix.",
        stacked=True,
    )

    flow_slide(
        prs, L,
        "A stay pair, concretely",
        [
            ("Two flights, same day, same city pair",
             "Same origin, destination, and aircraft type — but different routes"),
            ("One flight was delayed; its twin was not",
             "The twin's route is a real, demonstrated alternative — flown that same day"),
            ("The delayed flight kept its route anyway",
             "Its route is labelled preferred; the twin's route is the declined alternative"),
            ("This happened despite the cost",
             "Across all 6,672 such pairs: median peak delay 31 minutes, and 47.5% kept the pricier route in fuel"),
        ],
        "Real evidence that operators value stability itself — exactly what an over-eager model would violate.",
    )

    flow_slide(
        prs, L,
        "How the bias test works",
        [
            ("Find two similar flights, same day",
             "Same origin, destination, aircraft type — one delayed, one not"),
            ("The delayed flight had a reason to move",
             "Its twin's route was a real, demonstrated alternative — and it stayed anyway"),
            ("Label “stayed” as the preferred route",
             "Delay features are hidden, so route geometry alone must decide"),
            ("Score it without seeing delay",
             "Before this fix, the answer was 45.7% — worse than chance"),
        ],
        "This is the regulation-blind test behind the next chart — the fix raises this to 57.3%.",
    )

    image_slide(
        prs, L,
        "Correcting a bias toward change",
        "fig_staypairs.png",
        "With regulation clues hidden, route-only accuracy on stay "
        "pairs: 45.7% before correction, 57.3% after.",
    )

    flow_slide(
        prs, L,
        "Why the model sees differences, not totals",
        [
            ("Routes in different pairs vary hugely in scale",
             "A short regional hop and a long intercontinental sector burn very different absolute fuel"),
            ("So each indicator is rebased to its own pair",
             "Subtract the pair's minimum from both routes — the cheaper one becomes 0"),
            ("The model sees only the extra, not the total",
             "0 kg and 300 kg, not 2,000 kg and 2,300 kg — the gap that actually drove the choice"),
        ],
        "Skip this step, and the next chart shows exactly what it costs.",
    )

    image_slide(
        prs, L,
        "The one choice that changes everything",
        "fig_normalization.png",
        "2,000 kg alone is meaningless; the 300 kg gap is everything. "
        "Drop it: accuracy falls to chance (60.6% → 50.8%).",
    )

    # --- headline accuracy and where it lives ---

    flow_slide(
        prs, L,
        "How accuracy and confidence are measured",
        [
            ("Score both routes in a held-out pair", "Never used in training, tuning, or calibration"),
            ("Correct = higher score matches the choice", "Accuracy is this, averaged across all pairs: 68.3% overall"),
            ("The score gap is a confidence signal", "A bigger gap between the two scores means a more confident call"),
            ("Coverage = share of cases answered", "Answering only the surest fifth raises accuracy to 95%"),
        ],
        "The same measure drives both charts that follow.",
    )

    image_slide(
        prs, L,
        "The headline, and the curve behind it",
        "fig_coverage.png",
        "68.3% accuracy across all cases; 95% on the most confident fifth.",
    )

    image_slide(
        prs, L,
        "Where the accuracy lives",
        "fig_strata.png",
        "Accuracy rises with the delay being avoided — the larger the\n"
        "incentive, the clearer the decision.",
    )

    # --- what the model weighs ---

    flow_slide(
        prs, L,
        "How feature importance is measured",
        [
            ("Treat each input as a player in a game", "Route, fuel, delay, distance, and every other feature"),
            ("The payout is the model's score", "For one specific route, on one specific decision"),
            ("Add features one at a time, in random order", "Measure how much the score moves each time one joins"),
            ("Average that movement over the quarter", "Each feature's fair share of the credit — route structure earns the most"),
        ],
        "This ranks what moves the score, not what an operator consciously weighs.",
    )

    image_slide(
        prs, L,
        "What drives the model’s decisions",
        "fig_shap.png",
        "Route structure carries more weight than any single cost\n"
        "indicator — fuel, time, or distance.",
    )

    # --- from score to a real proposal, and its value ---

    flow_slide(
        prs, L,
        "How a score becomes a proposal",
        [
            ("The model scores each candidate route",
             "The alternative comes out +3.6 ahead — just the model's own number, not yet a percentage"),
            ("History says how sure to be",
             "Among past cases with a lead about this big, the airline picked the model's "
             "suggestion 8 times out of 10"),
            ("Only propose if that clears the bar",
             "The policy requires at least 6 of 10; 8 of 10 clears it, and it only counts as "
             "confirmed once the airline's later filing agrees"),
        ],
        "This is the policy behind every fuel figure in this talk.",
    )

    image_slide(
        prs, L,
        "The count behind that probability",
        "fig_calibration.png",
        "That's the whole mechanism — real past cases, counted, compared "
        "against the bar the policy requires.",
    )

    image_slide(
        prs, L,
        "From a score to a decision",
        "fig_tradeoff.png",
        "More confidence required means fewer suggestions but better precision. "
        "At this threshold: 16.4 kt fuel a quarter, ~200 kt CO₂/yr.",
    )

    flow_slide(
        prs, L,
        "A proposal, concretely",
        [
            ("Two candidate routes", "Same airline, aircraft, and departure time — no delay involved"),
            ("Comparison", "One route burns 2,016 kg less fuel — roughly a fifth of trip fuel"),
            ("Confidence check", "This case clears the threshold set for an 80%-precision channel"),
            ("Outcome", "The airline's later filing matched the suggested route exactly"),
        ],
        "One case from the test period, followed start to finish.",
    )

    image_slide(
        prs, L,
        "Where the value sits",
        "fig_concentration.png",
        "Six of 815 airlines account for half of the confirmed fuel;\n"
        "28 account for four-fifths.",
    )

    # --- does it survive time ---

    flow_slide(
        prs, L,
        "How the ageing test works",
        [
            ("Shift the whole design back one year", "Train through October 2025 instead of February 2026"),
            ("Freeze everything at the cutoff", "Same threshold, same calibration — no peeking at 2026 data"),
            ("Score month by month, unmodified", "January through June 2026, one month at a time"),
            ("Compare to the fresh model on the same months", "Isolates pure ageing — precision barely moves, as the next chart shows"),
        ],
        "The result: precision holds at 77-79% even six months out.",
    )

    image_slide(
        prs, L,
        "How fast does the model age?",
        "fig_horizon.png",
        "Precision holds steady (77-79%) six months on; it simply proposes "
        "less. Retraining recovers the 1-3 points that drift costs.",
    )

    limits_slide(
        prs, L,
        "Open questions and limitations",
        [
            "The model learns from what airlines did, not what they would say they prefer",
            "Only flights that changed are observed; those that stayed are not",
            "Every fuel figure is planned fuel, not fuel actually burned",
            "Cases where a better route was never proposed are invisible to this data",
            "A suggestion is only as reliable as the traffic picture it is computed against",
        ],
    )

    flow_slide(
        prs, L,
        "Deployed and serving",
        [
            ("Registered and version-controlled",
             "MLflow Model Registry, Unity Catalog, on MUAC's Databricks workspace"),
            ("Live REST endpoint",
             "Scores real candidate route sets on request; scales to zero when idle"),
            ("Ready for the next step",
             "A live shadow-mode trial is what would confirm real-world acceptance"),
        ],
        "The scoring system is live today — the next step is a trial in front of real dispatchers.",
    )

    two_card_slide(
        prs, L,
        "Where the full analysis is documented",
        [
            ("SESAR Innovation Days 2026", "Presented\nConference-length version", THEME_GREEN),
            ("Transportation Research\nPart C", "In preparation\nFull results and limitations", THEME_CYAN),
        ],
    )

    closing_slide(
        prs, L,
        "Questions?",
        "EUROCONTROL Innovation Hub  ·  Maastricht Upper Area Control Centre",
    )

    prs.save(OUT)
    print(f"wrote {OUT} ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
